from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn


OUT = Path('Apresentacao_Noivas_e_Cia.pptx')

# Palette aligned with the application design system.
BG = 'F6F7F9'
WHITE = 'FFFFFF'
INK = '111827'
MUTED = '475569'
BORDER = 'CBD5E1'
PINK = 'BE123C'
PINK_LIGHT = 'FFF1F3'
PINK_MID = 'FECDD6'
GREEN = '047857'
GREEN_LIGHT = 'ECFDF5'
AMBER = 'B45309'
AMBER_LIGHT = 'FFFBEB'
SLATE = '334155'
BLUE = '1D4ED8'
BLUE_LIGHT = 'EFF6FF'

# Shared layout grid — one left margin for every slide keeps the deck aligned.
MARGIN = 0.65
CRISP_RADIUS = 0.09  # standard corner radius for cards/panels (inches)


def color(hex_value):
    return RGBColor.from_string(hex_value)


def add_soft_shadow(shape, blur=0.06, dist=0.035, alpha=78):
    """Attach a subtle drop shadow so cards read as elevated, not flat rectangles."""
    sp_pr = shape.fill._xPr
    effect_lst = sp_pr.makeelement(qn('a:effectLst'), {})
    outer_shdw = sp_pr.makeelement(qn('a:outerShdw'), {
        'blurRad': str(Emu(Inches(blur))),
        'dist': str(Emu(Inches(dist))),
        'dir': '5400000',
        'rotWithShape': '0',
    })
    clr = sp_pr.makeelement(qn('a:srgbClr'), {'val': '1E293B'})
    alpha_el = sp_pr.makeelement(qn('a:alpha'), {'val': str(alpha * 1000)})
    clr.append(alpha_el)
    outer_shdw.append(clr)
    effect_lst.append(outer_shdw)
    sp_pr.append(effect_lst)


def add_shape(slide, shape_type, x, y, w, h, fill=WHITE, line=None, radius=False,
              corner_radius=CRISP_RADIUS, shadow=False):
    is_rounded = radius or shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else shape_type,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(line or fill)
    if is_rounded:
        min_dim = min(w, h)
        if min_dim > 0:
            shape.adjustments[0] = max(0.0, min(0.5, corner_radius / min_dim))
    if shadow:
        add_soft_shadow(shape)
    return shape


def add_text(slide, text, x, y, w, h, size=18, color_value=INK,
             bold=False, align=PP_ALIGN.LEFT, font='Aptos', valign=MSO_ANCHOR.TOP,
             line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color(color_value)
    p.space_after = Pt(0)
    if line_spacing:
        p.line_spacing = line_spacing
    return box


def bullet_list(slide, items, x, y, w, size=16, gap=0.44, bullet_color=PINK):
    for i, item in enumerate(items):
        yy = y + i * gap
        add_shape(slide, MSO_SHAPE.OVAL, x, yy + 0.1, 0.1, 0.1, fill=bullet_color)
        add_text(slide, item, x + 0.25, yy, w - 0.25, 0.33, size=size, color_value=SLATE)


def header(slide, title, subtitle=None, number=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.09, fill=PINK)
    add_text(slide, title, MARGIN, 0.42, 11.0, 0.5, size=30, bold=True)
    if subtitle:
        add_text(slide, subtitle, MARGIN, 0.99, 11.0, 0.34, size=13, color_value=MUTED)
    if number is not None:
        add_text(slide, f'{number:02d}', 12.05, 0.48, 0.65, 0.3, size=11,
                 color_value=PINK, bold=True, align=PP_ALIGN.RIGHT)


def footer(slide):
    add_text(slide, 'NOIVAS & CIA  •  APRESENTAÇÃO DO SISTEMA', MARGIN, 7.08, 5.1, 0.2,
             size=8, color_value=MUTED, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 11.97, 7.13, 0.72, 0.04, fill=PINK)


def card(slide, x, y, w, h, title, body, accent=PINK, body_size=14):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=WHITE, line=BORDER, shadow=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.02, y + 0.1, 0.07, h - 0.2, fill=accent)
    add_text(slide, title, x + 0.28, y + 0.25, w - 0.5, 0.35, size=16, bold=True)
    add_text(slide, body, x + 0.28, y + 0.74, w - 0.55, h - 0.9, size=body_size,
             color_value=MUTED, line_spacing=1.15)


def tag(slide, text, x, y, w, fill=PINK_LIGHT, text_color=PINK):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.32, fill=fill, line=fill,
              corner_radius=0.16)
    add_text(slide, text, x + 0.06, y + 0.075, w - 0.12, 0.17, size=9,
             color_value=text_color, bold=True, align=PP_ALIGN.CENTER)


def step(slide, x, y, title, detail, number, accent=PINK):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.85, 1.42, fill=WHITE, line=BORDER, shadow=True)
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.16, y + 0.16, 0.38, 0.38, fill=accent)
    add_text(slide, str(number), x + 0.16, y + 0.235, 0.38, 0.15, size=10,
             color_value=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.16, y + 0.66, 1.55, 0.22, size=13, bold=True)
    add_text(slide, detail, x + 0.16, y + 0.95, 1.53, 0.4, size=10, color_value=MUTED,
             line_spacing=1.1)


def arrow(slide, x, y):
    add_text(slide, '›', x, y - 0.02, 0.28, 0.45, size=31, bold=False, color_value=PINK,
             align=PP_ALIGN.CENTER)


def make_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = 'Noivas & Cia — Apresentação do Sistema'
    prs.core_properties.subject = 'Evolução do sistema local para uma aplicação web'
    prs.core_properties.author = 'Noivas & Cia'
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=BG)
    add_shape(slide, MSO_SHAPE.OVAL, 9.5, -1.7, 5.2, 5.2, fill=PINK_LIGHT)
    add_shape(slide, MSO_SHAPE.OVAL, 10.75, -0.6, 3.6, 3.6, fill=PINK_MID)
    add_shape(slide, MSO_SHAPE.OVAL, 9.18, 3.75, 4.9, 4.9, fill=WHITE)
    add_text(slide, 'NOIVAS & CIA', MARGIN, 0.78, 3.1, 0.28, size=12, color_value=PINK, bold=True)
    add_text(slide, 'Uma nova forma de\ngerenciar locações.', MARGIN, 1.46, 7.8, 1.5,
             size=42, bold=True, line_spacing=1.02)
    add_text(slide, 'Do sistema local para uma experiência web mais simples,\norganizada e preparada para o futuro.',
             MARGIN, 3.3, 6.7, 0.8, size=18, color_value=MUTED, line_spacing=1.2)
    tag(slide, 'APRESENTAÇÃO PARA CLIENTE', MARGIN, 4.4, 2.15)
    # Decorative web-app card.
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.38, 2.0, 3.75, 3.52, fill=WHITE, line=WHITE, shadow=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.38, 2.0, 3.75, 0.38, fill=PINK)
    for xx in (8.62, 8.82, 9.02):
        add_shape(slide, MSO_SHAPE.OVAL, xx, 2.12, 0.08, 0.08, fill=WHITE)
    add_text(slide, 'Painel do dia', 8.72, 2.72, 2.5, 0.28, size=16, bold=True)
    for xx, label, value, fill, text_color in [
        (8.7, 'A retirar', '08', PINK_LIGHT, PINK),
        (10.37, 'Em aberto', '12', AMBER_LIGHT, AMBER),
    ]:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, xx, 3.2, 1.38, 0.86, fill=fill, line=fill)
        add_text(slide, value, xx + 0.12, 3.34, 0.4, 0.25, size=20, bold=True, color_value=text_color)
        add_text(slide, label, xx + 0.12, 3.71, 1.05, 0.15, size=8, color_value=MUTED)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 8.7, 4.36, 2.7, 0.63, fill=GREEN_LIGHT, line=GREEN_LIGHT)
    add_text(slide, '✓  Acervo sob controle', 8.92, 4.58, 2.2, 0.18, size=11, bold=True, color_value=GREEN)
    add_text(slide, 'Julho de 2026', MARGIN, 6.65, 2.2, 0.22, size=11, color_value=MUTED)

    # 2. Challenge
    slide = prs.slides.add_slide(blank)
    header(slide, 'Por que evoluir agora?', 'A operação merece uma ferramenta compatível com o ritmo da loja.', 2)
    add_text(slide, 'O sistema anterior cumpre sua história, mas limita o presente.', MARGIN, 1.55, 8.8, 0.35,
             size=20, bold=True)
    card(slide, 0.65, 2.18, 3.83, 2.6, 'Dependência de um único ponto',
         'A operação fica presa a máquinas e a um ambiente antigo, aumentando o risco de parada.', AMBER)
    card(slide, 4.75, 2.18, 3.83, 2.6, 'Informação menos acessível',
         'Encontrar dados de clientes, peças e contratos exige mais tempo durante o atendimento.', PINK)
    card(slide, 8.85, 2.18, 3.83, 2.6, 'Visão limitada do dia',
         'Pendências, devoluções e recebimentos precisam ser conferidos com mais esforço.', BLUE)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.65, 5.35, 12.02, 0.83, fill=PINK_LIGHT, line=PINK_LIGHT)
    add_text(slide, 'Objetivo da modernização: reduzir riscos, ganhar agilidade e ter tudo o que importa visível na hora certa.',
             0.98, 5.61, 11.3, 0.24, size=16, bold=True, color_value=PINK, align=PP_ALIGN.CENTER)
    footer(slide)

    # 3. Web transition
    slide = prs.slides.add_slide(blank)
    header(slide, 'Do aplicativo local para o sistema web', 'A mudança é percebida no dia a dia — sem complicar a rotina.', 3)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.65, 1.65, 5.65, 4.62, fill=WHITE, line=BORDER, shadow=True)
    tag(slide, 'ANTES', 0.97, 1.98, 0.9, fill=AMBER_LIGHT, text_color=AMBER)
    add_text(slide, 'Sistema instalado', 0.98, 2.53, 3.7, 0.34, size=23, bold=True)
    bullet_list(slide, [
        'Uso condicionado a uma máquina específica.',
        'Tecnologia antiga, mais difícil de manter.',
        'Menos praticidade para consultar a operação.',
    ], 1.0, 3.25, 4.85, size=15, gap=0.62, bullet_color=AMBER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.02, 1.65, 5.65, 4.62, fill=PINK_LIGHT, line=PINK_MID, shadow=True)
    tag(slide, 'AGORA', 7.34, 1.98, 0.9)
    add_text(slide, 'Sistema acessado no navegador', 7.35, 2.53, 4.7, 0.34, size=23, bold=True, color_value=PINK)
    bullet_list(slide, [
        'Acesso por computadores autorizados na rede.',
        'Interface moderna, clara e responsiva.',
        'Dados centralizados em um único sistema.',
    ], 7.37, 3.25, 4.85, size=15, gap=0.62)
    add_text(slide, 'A mesma operação, com mais liberdade, clareza e controle.', 0.65, 6.57, 12.0, 0.25,
             size=16, bold=True, color_value=MUTED, align=PP_ALIGN.CENTER)
    footer(slide)

    # 4. Flow
    slide = prs.slides.add_slide(blank)
    header(slide, 'A rotina da loja em um fluxo simples', 'Cada etapa conversa com a próxima, evitando controles paralelos.', 4)
    step(slide, 0.58, 2.02, 'Cliente', 'Cadastro e histórico sempre à mão.', 1, BLUE)
    arrow(slide, 2.53, 2.47)
    step(slide, 2.84, 2.02, 'Acervo', 'Consulta de peça e disponibilidade.', 2, PINK)
    arrow(slide, 4.79, 2.47)
    step(slide, 5.10, 2.02, 'Locação', 'Contrato, itens e datas organizados.', 3, AMBER)
    arrow(slide, 7.05, 2.47)
    step(slide, 7.36, 2.02, 'Retirada e retorno', 'Acompanhamento do que saiu e voltou.', 4, GREEN)
    arrow(slide, 9.31, 2.47)
    step(slide, 9.62, 2.02, 'Financeiro', 'Parcelas, pagamentos e pendências.', 5, BLUE)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.28, 4.65, 10.72, 0.94, fill=WHITE, line=BORDER, shadow=True)
    add_text(slide, 'Resultado: atendimento mais fluido, menos consultas manuais e decisões com informação atualizada.',
             1.63, 4.96, 10.0, 0.25, size=17, bold=True, color_value=SLATE, align=PP_ALIGN.CENTER)
    footer(slide)

    # 5. Improvements / operational dashboard
    slide = prs.slides.add_slide(blank)
    header(slide, 'Melhorias que aparecem na operação', 'Mais informação útil na tela, no momento em que a equipe precisa.', 5)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.65, 1.6, 6.05, 4.95, fill=WHITE, line=BORDER, shadow=True)
    add_text(slide, 'Painel de acompanhamento', 0.96, 1.95, 3.5, 0.28, size=18, bold=True)
    add_text(slide, 'Resumo da rotina com prioridades visíveis.', 0.96, 2.32, 4.4, 0.22, size=12, color_value=MUTED)
    for x, title, value, fill, text_color in [
        (0.96, 'A retirar', '08', PINK_LIGHT, PINK),
        (2.78, 'A devolver', '05', AMBER_LIGHT, AMBER),
        (4.60, 'Em aberto', '12', BLUE_LIGHT, BLUE),
    ]:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.85, 1.48, 1.0, fill=fill, line=fill)
        add_text(slide, value, x + 0.14, 3.05, 0.5, 0.28, size=22, bold=True, color_value=text_color)
        add_text(slide, title, x + 0.14, 3.50, 1.18, 0.14, size=9, color_value=MUTED)
    add_text(slide, 'Próximas retiradas', 0.96, 4.35, 2.0, 0.2, size=13, bold=True)
    for yy, name, status, fill in [(4.75, 'Contrato #1034  •  Ana Souza', 'Hoje', PINK_LIGHT),
                                   (5.22, 'Contrato #1037  •  Beatriz Lima', 'Amanhã', BLUE_LIGHT)]:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.96, yy, 5.25, 0.36, fill=fill, line=fill)
        add_text(slide, name, 1.13, yy + 0.10, 3.45, 0.15, size=10, color_value=SLATE)
        add_text(slide, status, 5.15, yy + 0.10, 0.7, 0.15, size=10, bold=True, color_value=PINK, align=PP_ALIGN.RIGHT)
    card(slide, 7.24, 1.6, 5.43, 1.2, 'Disponibilidade confiável',
         'Antes de confirmar, a equipe consulta se a peça estará livre na data desejada.', GREEN, body_size=13)
    card(slide, 7.24, 3.05, 5.43, 1.2, 'Status de cada locação',
         'Acompanhe o que está reservado, retirado, devolvido ou em atraso.', PINK, body_size=13)
    card(slide, 7.24, 4.5, 5.43, 1.2, 'Histórico preservado',
         'Clientes e contratos ficam organizados para consultas futuras.', BLUE, body_size=13)
    footer(slide)

    # 6. Finance
    slide = prs.slides.add_slide(blank)
    header(slide, 'Financeiro com mais clareza e acompanhamento', 'O sistema transforma pendências em informação prática para a equipe.', 6)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.65, 1.72, 4.1, 4.5, fill=PINK_LIGHT, line=PINK_MID, shadow=True)
    add_text(slide, 'Controle financeiro', 1.02, 2.05, 3.2, 0.3, size=22, bold=True, color_value=PINK)
    add_text(slide, 'Cada locação pode ter suas parcelas, pagamentos e saldo acompanhados de perto.',
             1.02, 2.58, 3.1, 0.82, size=16, color_value=SLATE, line_spacing=1.2)
    tag(slide, 'TRANSPARÊNCIA PARA DECIDIR', 1.02, 4.22, 2.45, fill=WHITE, text_color=PINK)
    add_text(slide, '• Valores e vencimentos\n• Pagamentos parciais\n• Juros configurados para atrasos\n• Visão de contas em aberto',
             1.02, 4.8, 3.1, 1.04, size=13, color_value=SLATE)
    card(slide, 5.22, 1.72, 3.35, 2.0, 'Para o balcão',
         'Registro de pagamento simples e saldo atualizado na hora.', GREEN, body_size=14)
    card(slide, 8.95, 1.72, 3.35, 2.0, 'Para a gestão',
         'Relatórios ajudam a enxergar valores a receber e vencimentos.', BLUE, body_size=14)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.22, 4.18, 7.08, 1.3, fill=WHITE, line=BORDER, shadow=True)
    add_text(slide, '“O financeiro deixa de ser apenas registro e passa a apoiar a decisão diária.”',
             5.7, 4.58, 6.1, 0.35, size=17, bold=True, color_value=SLATE, align=PP_ALIGN.CENTER)
    footer(slide)

    # 7. Data and security
    slide = prs.slides.add_slide(blank)
    header(slide, 'Informações organizadas, acesso controlado', 'Tudo fica conectado sem perder a simplicidade de uso.', 7)
    add_shape(slide, MSO_SHAPE.OVAL, 5.55, 2.31, 2.25, 2.25, fill=PINK, line=PINK)
    add_text(slide, 'Dados da\nNoivas & Cia', 5.76, 3.02, 1.85, 0.62, size=18, bold=True,
             color_value=WHITE, align=PP_ALIGN.CENTER)
    for x, y, title, detail, accent in [
        (0.78, 1.7, 'Clientes', 'Cadastros e histórico.', BLUE),
        (0.78, 4.75, 'Acervo', 'Peças, categorias e disponibilidade.', GREEN),
        (9.15, 1.7, 'Locações', 'Contratos, datas e movimentações.', AMBER),
        (9.15, 4.75, 'Financeiro', 'Parcelas, pagamentos e relatórios.', PINK),
    ]:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 3.3, 1.15, fill=WHITE, line=BORDER, shadow=True)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.02, y + 0.1, 0.06, 0.95, fill=accent)
        add_text(slide, title, x + 0.25, y + 0.24, 2.55, 0.22, size=15, bold=True)
        add_text(slide, detail, x + 0.25, y + 0.62, 2.7, 0.21, size=11, color_value=MUTED)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.18, 5.94, 7.0, 0.58, fill=GREEN_LIGHT, line=GREEN_LIGHT)
    add_text(slide, 'Usuários têm acesso somente ao que precisam; a gestão mantém o controle.',
             3.45, 6.14, 6.5, 0.17, size=12, bold=True, color_value=GREEN, align=PP_ALIGN.CENTER)
    footer(slide)

    # 8. Reliability & continuation
    slide = prs.slides.add_slide(blank)
    header(slide, 'Mais segurança para seguir em frente', 'A modernização reduz dependências e cria uma base mais confiável para a operação.', 8)
    card(slide, 0.65, 1.75, 3.82, 3.55, 'Acesso por usuário',
         'Cada pessoa entra com sua própria conta. As áreas do sistema podem ser liberadas conforme a função.', PINK, body_size=15)
    card(slide, 4.76, 1.75, 3.82, 3.55, 'Cópias de segurança',
         'O sistema possui rotina de backup para preservar a base de dados e os arquivos importantes.', GREEN, body_size=15)
    card(slide, 8.87, 1.75, 3.82, 3.55, 'Base preparada para evoluir',
         'A estrutura foi organizada por áreas da loja, tornando futuras melhorias mais seguras e graduais.', BLUE, body_size=15)
    add_text(slide, 'Modernizar não é só trocar de tela: é proteger a continuidade do negócio.', 1.18, 6.03, 10.9, 0.3,
             size=20, bold=True, color_value=SLATE, align=PP_ALIGN.CENTER)
    footer(slide)

    # 9. Closing
    slide = prs.slides.add_slide(blank)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, fill=PINK)
    add_shape(slide, MSO_SHAPE.OVAL, -1.0, 5.2, 4.2, 4.2, fill='9F1239')
    add_shape(slide, MSO_SHAPE.OVAL, 10.5, -1.5, 4.2, 4.2, fill='9F1239')
    add_text(slide, 'NOIVAS & CIA', 0.75, 0.8, 2.8, 0.27, size=12, color_value=PINK_MID, bold=True)
    add_text(slide, 'Mais controle para a equipe.\nMais confiança para o negócio.', 0.75, 1.62, 10.0, 1.25,
             size=35, bold=True, color_value=WHITE)
    add_text(slide, 'O novo sistema reúne a operação em uma experiência clara, moderna e preparada\npara acompanhar o crescimento da Noivas & Cia.',
             0.77, 3.33, 9.6, 0.72, size=18, color_value='FFE4E9')
    for x, text in [(0.78, 'Atendimento mais ágil'), (4.4, 'Dados organizados'), (7.62, 'Gestão mais segura')]:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 5.02, 2.84, 0.56, fill=WHITE, line=WHITE)
        add_text(slide, text, x + 0.08, 5.22, 2.68, 0.16, size=11, bold=True, color_value=PINK,
                 align=PP_ALIGN.CENTER)
    add_text(slide, 'Agradecemos!', 0.78, 6.54, 2.0, 0.28, size=16, bold=True, color_value=WHITE)

    prs.save(OUT)
    return OUT


if __name__ == '__main__':
    print(make_presentation())
